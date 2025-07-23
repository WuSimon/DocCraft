import os
import argparse
import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from io import BytesIO
from PIL import Image
import re

def html_to_pptx(input_html, output_pptx, download_images=True):
    """
    Convert an HTML file to a PowerPoint presentation
    
    Args:
        input_html (str): Path to input HTML file
        output_pptx (str): Path for output PPTX file
        download_images (bool): Whether to download remote images
    """
    # Read HTML content
    with open(input_html, 'r', encoding='utf-8') as f:
        html_content = f.read()
    
    # Parse HTML with BeautifulSoup
    soup = BeautifulSoup(html_content, 'html.parser')
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)
    
    # Add slide
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    bg_color = extract_background_color(soup)
    fill.fore_color.rgb = bg_color or RGBColor(30, 30, 30)
    
    # Extract and add title
    title = soup.find('h1', class_='title')
    if title:
        add_text_element(slide, title.text, 
                         left=0.5, top=0.3, width=12, height=0.8,
                         font_size=36, color=RGBColor(74, 134, 232), bold=True)
    
    # Extract and add header
    header = soup.find(class_='header')
    if header:
        add_text_element(slide, header.text.strip(), 
                         left=10, top=0.5, width=2.5, height=0.4,
                         font_size=12, color=RGBColor(160, 174, 192))
    
    # Extract and add source link
    source_link = soup.find(class_='source-link')
    if source_link:
        add_text_element(slide, source_link.text.strip(), 
                         left=0.5, top=1.4, width=6, height=0.4,
                         font_size=14, color=RGBColor(74, 134, 232))
    
    # Extract and add command
    command = soup.find(class_='command')
    if command:
        add_text_element(slide, command.text.strip(), 
                         left=0.5, top=1.8, width=6, height=0.8,
                         font_size=14, color=RGBColor(224, 224, 224), 
                         font_name='Courier New')
    
    # Extract and add image
    img = soup.find('img')
    if img and 'src' in img.attrs:
        img_url = img['src']
        try:
            if download_images:
                response = requests.get(img_url)
                img_bytes = BytesIO(response.content)
                slide.shapes.add_picture(
                    img_bytes, 
                    left=Inches(0.5), 
                    top=Inches(2.7), 
                    width=Inches(6)
                )
            else:
                add_placeholder(slide, "Remote Image", 
                                left=0.5, top=2.7, width=6, height=4)
        except Exception as e:
            add_placeholder(slide, "Image Load Failed", 
                            left=0.5, top=2.7, width=6, height=4)
    
    # Extract output boxes
    output_boxes = soup.find_all(class_='output-box')
    for i, box in enumerate(output_boxes):
        # Extract box title
        title_elem = box.find(class_='parser-title')
        title_text = title_elem.text.strip() if title_elem else f"Output Box {i+1}"
        
        # Extract content
        content_elem = box.find(class_='json-content')
        content_text = content_elem.text.strip() if content_elem else ""
        
        # Determine color based on box type
        if 'qwenvl' in box.get('class', []):
            color = RGBColor(52, 168, 83)  # Green
        elif 'deepseekvl' in box.get('class', []):
            color = RGBColor(234, 67, 53)   # Red
        elif 'layoutlmv3' in box.get('class', []):
            color = RGBColor(251, 188, 5)   # Yellow
        else:
            color = RGBColor(160, 174, 192) # Default
        
        # Add output box
        top = 1.8 + i * 1.8
        add_text_element(slide, title_text, 
                         left=7, top=top, width=5.5, height=0.4,
                         font_size=16, color=color, bold=True)
        
        add_text_element(slide, content_text, 
                         left=7, top=top + 0.4, width=5.5, height=1.2,
                         font_size=11, color=RGBColor(200, 200, 200), 
                         font_name='Courier New')
    
    # Extract and add footer
    footer = soup.find(class_='footer')
    if footer:
        footer_text = ' '.join(footer.stripped_strings)
        add_text_element(slide, footer_text, 
                         left=0.5, top=6.8, width=12, height=0.4,
                         font_size=14, color=RGBColor(160, 174, 192))
    
    # Save presentation
    prs.save(output_pptx)
    print(f"Successfully created PowerPoint presentation: {output_pptx}")

def add_text_element(slide, text, left, top, width, height, 
                     font_size=12, color=RGBColor(0, 0, 0), 
                     bold=False, font_name=None, alignment=PP_ALIGN.LEFT):
    """
    Add a text element to the slide
    
    Args:
        slide: PowerPoint slide object
        text: Text content
        left, top, width, height: Position and size in inches
        font_size: Font size in points
        color: Font color
        bold: Whether text is bold
        font_name: Font family name
        alignment: Text alignment
    """
    textbox = slide.shapes.add_textbox(
        left=Inches(left), 
        top=Inches(top), 
        width=Inches(width), 
        height=Inches(height))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    
    # Clear default empty paragraph
    for paragraph in text_frame.paragraphs:
        p = paragraph._element
        p.getparent().remove(p)
    
    # Add new paragraph with text
    p = text_frame.add_paragraph()
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    if font_name:
        run.font.name = font_name

def add_placeholder(slide, text, left, top, width, height):
    """
    Add a placeholder shape with text
    
    Args:
        slide: PowerPoint slide object
        text: Placeholder text
        left, top, width, height: Position and size in inches
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(left), Inches(top), 
        Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(45, 45, 45)
    shape.line.color.rgb = RGBColor(100, 100, 100)
    
    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(180, 180, 180)

def extract_background_color(soup):
    """
    Extract background color from HTML
    
    Args:
        soup: BeautifulSoup object
        
    Returns:
        RGBColor object or None
    """
    try:
        # Look for background color in style attributes
        for style in soup.find_all(style=True):
            style_text = style['style']
            if 'background:' in style_text:
                match = re.search(r'background:\s*#([0-9a-fA-F]{6})', style_text)
                if match:
                    hex_color = match.group(1)
                    return RGBColor.from_string(f'#{hex_color}')
        
        # Look for background-color specifically
        for style in soup.find_all(style=True):
            style_text = style['style']
            if 'background-color:' in style_text:
                match = re.search(r'background-color:\s*#([0-9a-fA-F]{6})', style_text)
                if match:
                    hex_color = match.group(1)
                    return RGBColor.from_string(f'#{hex_color}')
    except:
        pass
    
    return None

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description='Convert HTML to PowerPoint')
    parser.add_argument('input_html', help='Input HTML file path')
    parser.add_argument('output_pptx', help='Output PPTX file path')
    parser.add_argument('--no-images', action='store_true', 
                        help='Skip downloading remote images')
    
    args = parser.parse_args()
    
    # Ensure output directory exists
    output_dir = os.path.dirname(args.output_pptx)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    html_to_pptx(
        input_html=args.input_html,
        output_pptx=args.output_pptx,
        download_images=not args.no_images
    )