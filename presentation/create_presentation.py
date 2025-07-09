from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime

def add_flowchart(slide, left, top, width, height):
    """Add a simple flowchart to the slide."""
    # Document box
    doc_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width/3, height/3)
    doc_box.text = "Document"
    doc_box.fill.solid()
    doc_box.fill.fore_color.rgb = RGBColor(0, 112, 192)
    
    # Arrow
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + width/3, top + height/6, width/6, height/6)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(0, 112, 192)
    
    # DocCraft box
    craft_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + width/2, top, width/3, height/3)
    craft_box.text = "DocCraft"
    craft_box.fill.solid()
    craft_box.fill.fore_color.rgb = RGBColor(0, 176, 80)
    
    # Arrow
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + 5*width/6, top + height/6, width/6, height/6)
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = RGBColor(0, 176, 80)
    
    # Data box
    data_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + width - width/3, top, width/3, height/3)
    data_box.text = "Structured Data"
    data_box.fill.solid()
    data_box.fill.fore_color.rgb = RGBColor(255, 192, 0)

def create_title_slide(prs):
    """Create the title slide with flowchart."""
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "DocCraft: Streamlining Document Intelligence for Smarter Decisions"
    subtitle.text = f"A Python Package for Unified Document Parsing & Benchmarking\nSimon\n{datetime.now().strftime('%B %d, %Y')}"
    
    # Style the title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Add flowchart below the subtitle
    add_flowchart(slide, Inches(2), Inches(4.5), Inches(12), Inches(2))

def create_problem_slide(prs):
    """Create the problem slide with visual elements."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "The Problem (Gastronomy Case Study)"
    content.text = "Hook: \"Why are restaurant owners struggling to track supplier costs?\"\n\n" \
                  "• Friends in gastronomy face price opacity\n" \
                  "• Manual invoice processing is error-prone and time-consuming\n" \
                  "• Key Insight: Automating this requires parsing unstructured documents with high accuracy"
    
    # Add a stylized invoice shape on the right side
    invoice = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9), Inches(2.5), Inches(3), Inches(4))
    invoice.text = "📄 Invoice\n\n$ Cost\n📦 Items\n📋 Categories"
    invoice.fill.solid()
    invoice.fill.fore_color.rgb = RGBColor(240, 240, 240)

def create_motivation_slide(prs):
    """Create the motivation slide with puzzle pieces."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "The Motivation (Your \"Aha!\" Moment)"
    content.text = "Challenge: \"Why can't existing tools solve this?\"\n\n" \
                  "• Fragmented ecosystem: OCR, PDF parsers, AI models work in silos\n" \
                  "• No easy way to compare performance across models/formats\n" \
                  "• Your Goal: Build a unified tool to solve your friends' problem and empower developers"
    
    # Add puzzle pieces in a row at the bottom
    for i, label in enumerate(["OCR", "PDF", "AI Models"]):
        piece = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4 + i*2.5), Inches(5.5), Inches(2), Inches(1))
        piece.text = label
        piece.fill.solid()
        piece.fill.fore_color.rgb = RGBColor(0, 112, 192)

def create_intro_slide(prs):
    """Create the introduction slide with pipeline diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Introducing DocCraft"
    content.text = "Tagline: \"One Package to Parse Them All.\"\n\n" \
                  "Core Features:\n" \
                  "• Unified interface for OCR (Tesseract, AWS Textract), PDF (PyPDF, Camelot), AI models\n" \
                  "• Benchmarking suite: Compare speed/accuracy across tools\n" \
                  "• Customizable pipelines (e.g., OCR → AI model → CSV output)"
    
    # Add pipeline diagram below the content
    add_flowchart(slide, Inches(2), Inches(5), Inches(12), Inches(2))

def create_architecture_slide1(prs):
    """Create the first architecture slide with layered diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Technical Architecture (1/2)"
    content.text = "How DocCraft Works Under the Hood\n\n" \
                  "Components:\n" \
                  "• Input Adapters: PDFs, images, scans\n" \
                  "• Parsing Layer: OCR/PDF extraction + AI model integration\n" \
                  "• Benchmarking Engine: Metrics (F1 score, latency) across pipelines"
    
    # Add layered architecture diagram on the right
    layers = ["Input Adapters", "Parsing Layer", "Benchmarking Engine"]
    colors = [RGBColor(0, 112, 192), RGBColor(0, 176, 80), RGBColor(255, 192, 0)]
    
    for i, (layer, color) in enumerate(zip(layers, colors)):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9), Inches(2.5 + i*1.5), Inches(3), Inches(1))
        shape.text = layer
        shape.fill.solid()
        shape.fill.fore_color.rgb = color

def create_architecture_slide2(prs):
    """Create the second architecture slide with demo flow."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Technical Architecture (2/2) – Use Case Demo"
    content.text = "Demo Flow: Restaurant Invoice → DocCraft → Structured Data\n\n" \
                  "Step 1: Extract text (OCR/PDF)\n" \
                  "Step 2: AI model identifies line items, costs, categories\n" \
                  "Step 3: Output to CSV/database + benchmarking report"
    
    # Add side-by-side comparison below the content
    invoice = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(5), Inches(3), Inches(2))
    invoice.text = "📄 Invoice\n\n$ Cost\n📦 Items\n📋 Categories"
    invoice.fill.solid()
    invoice.fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.5), Inches(5.5), Inches(1), Inches(1))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(0, 112, 192)
    
    table = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8), Inches(5), Inches(3), Inches(2))
    table.text = "📊 Table\n\nCost | Items\nCategories"
    table.fill.solid()
    table.fill.fore_color.rgb = RGBColor(0, 176, 80)

def create_impact_slide(prs):
    """Create the impact slide with icons."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Why This Matters"
    content.text = "Impact for Gastronomy:\n" \
                  "• Transparent pricing databases → better supplier negotiations\n" \
                  "• Saves hours/week on manual data entry\n\n" \
                  "Broader Applications:\n" \
                  "• Legal contracts, medical records, academic research"
    
    # Add industry icons in a row at the bottom
    icons = ["👨‍🍳", "⚖️", "🏥", "📚"]
    for i, icon in enumerate(icons):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3 + i*2.5), Inches(5.5), Inches(1.5), Inches(1.5))
        shape.text = icon
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0, 112, 192)

def create_roadmap_slide(prs):
    """Create the roadmap slide with GitHub logo."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Roadmap & Call to Action"
    content.text = "Next Steps:\n" \
                  "• Open-source MVP in 3 months\n" \
                  "• Expand to handwritten text, multilingual support\n\n" \
                  "Ask: \"Interested in collaborating? Let's build DocCraft together!\""
    
    # Add GitHub logo and banner on the right
    github = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9), Inches(2.5), Inches(3), Inches(3))
    github.text = "🐙 GitHub"
    github.fill.solid()
    github.fill.fore_color.rgb = RGBColor(36, 41, 46)
    
    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9), Inches(6), Inches(3), Inches(0.5))
    banner.text = "Coming Soon!"
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(255, 192, 0)

def create_thank_you_slide(prs):
    """Create the thank you slide with quote box."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    title.text = "Thank You!"
    
    # Add quote box centered
    left = Inches(3)
    top = Inches(3)
    width = Inches(10)
    height = Inches(2)
    
    quote_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    quote_box.text = "Closing Quote: \"Automating the mundane to focus on what matters.\""
    quote_box.fill.solid()
    quote_box.fill.fore_color.rgb = RGBColor(0, 112, 192)
    
    qa_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top + height + Inches(0.5), width, height/2)
    qa_box.text = "Q&A Prompt: \"How would YOU use intelligent document parsing?\""
    qa_box.fill.solid()
    qa_box.fill.fore_color.rgb = RGBColor(0, 176, 80)

def create_presentation():
    """Create the complete presentation."""
    prs = Presentation()
    
    # Set slide dimensions to 16:9
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Create slides
    create_title_slide(prs)
    create_problem_slide(prs)
    create_motivation_slide(prs)
    create_intro_slide(prs)
    create_architecture_slide1(prs)
    create_architecture_slide2(prs)
    create_impact_slide(prs)
    create_roadmap_slide(prs)
    create_thank_you_slide(prs)
    
    # Save the presentation
    prs.save('DocCraft_Presentation_Draft.pptx')

if __name__ == "__main__":
    create_presentation() 